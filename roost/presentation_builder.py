"""Presentation builder — python-pptx wrapper for creating .pptx files.

Builds PowerPoint presentations from structured content dicts (as produced
by meeting_notes_service.generate_presentation_content).
"""

import logging
import os

logger = logging.getLogger("roost.presentation_builder")


def build_presentation_from_content(
    content: dict,
    template_path: str | None = None,
    output_path: str = "/tmp/presentation.pptx",
) -> str:
    """Build a .pptx file from structured content.

    Args:
        content: Dict with keys: title, subtitle, slides.
            Each slide has: title, bullets (list[str]), notes (str).
        template_path: Optional .pptx template file. If provided, uses its
            slide layouts. Otherwise creates a blank presentation.
        output_path: Where to save the .pptx file.

    Returns:
        The output file path.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    if template_path and os.path.isfile(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Find layouts by name for template compatibility, fall back to index
    layouts = prs.slide_layouts
    title_layout = _find_layout(layouts, ["Title Slide", "title"], fallback_idx=0)
    content_layout = _find_layout(
        layouts,
        ["Title and Content", "Title, Content", "content"],
        fallback_idx=1,
    )

    pres_title = content.get("title", "Presentation")
    pres_subtitle = content.get("subtitle", "")
    slides_data = content.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        slide_title = slide_data.get("title", f"Slide {i + 1}")
        bullets = slide_data.get("bullets", [])
        notes_text = slide_data.get("notes", "")

        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(title_layout)
            _set_placeholder_text(slide, 0, pres_title)
            _set_placeholder_text(slide, 1, pres_subtitle or slide_title)
        else:
            # Content slide
            slide = prs.slides.add_slide(content_layout)
            _set_placeholder_text(slide, 0, slide_title)

            if bullets:
                body_ph = _get_placeholder(slide, 1)
                if body_ph and body_ph.has_text_frame:
                    tf = body_ph.text_frame
                    tf.clear()
                    for j, bullet in enumerate(bullets):
                        if j == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    # Style bullets
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(18)
                else:
                    # No body placeholder — add a text box
                    from pptx.util import Emu
                    left = Inches(0.8)
                    top = Inches(1.8)
                    width = Inches(8.4)
                    height = Inches(4.5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for j, bullet in enumerate(bullets):
                        if j == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                        tf.paragraphs[-1].font.size = Pt(18)

        # Add speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    size = os.path.getsize(output_path)
    logger.info(
        "Built presentation '%s': %d slides, %d bytes → %s",
        pres_title, len(slides_data), size, output_path,
    )
    return output_path


def build_presentation_from_template(
    template_path: str,
    replacements: dict[str, str],
    output_path: str,
) -> dict:
    """Fill a template .pptx with placeholder replacements.

    Reuses the existing slides_service.pptx_replace_text() function.

    Args:
        template_path: Path to the template .pptx file.
        replacements: Dict mapping placeholder strings to values.
        output_path: Where to save the filled presentation.

    Returns:
        Dict with output_path, total_replacements, per_placeholder counts.
    """
    import shutil
    from roost.slides_service import pptx_replace_text

    # Copy template to output location first
    shutil.copy2(template_path, output_path)

    # Apply replacements on the copy
    return pptx_replace_text(output_path, replacements)


# ── Helpers ──────────────────────────────────────────────────────────


def _find_layout(layouts, names: list[str], fallback_idx: int = 0):
    """Find a slide layout by name (case-insensitive), fall back to index."""
    for layout in layouts:
        for name in names:
            if name.lower() in layout.name.lower():
                return layout
    # Fallback to index if available
    if fallback_idx < len(layouts):
        return layouts[fallback_idx]
    return layouts[0]


def _get_placeholder(slide, idx: int):
    """Get a placeholder by index, or None if not found."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def _set_placeholder_text(slide, idx: int, text: str):
    """Set text on a placeholder by index, silently skip if not found."""
    ph = _get_placeholder(slide, idx)
    if ph and ph.has_text_frame:
        ph.text_frame.text = text
