"""Slide editing service — PowerPoint (python-pptx) + Google Slides API.

Shared service layer used by the MCP tools module.
"""

import logging
import os

logger = logging.getLogger("roost.slides")


# ── PowerPoint (python-pptx) ─────────────────────────────────────────


def pptx_list_placeholders(file_path: str) -> list[dict]:
    """List all slides with their shapes/placeholders and text content.

    Returns a list of dicts, one per slide, each containing:
      - slide_number (1-based)
      - shapes: list of {name, shape_id, text, placeholder_idx}
    """
    from pptx import Presentation

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    result = []

    for i, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "shape_id": shape.shape_id,
                "text": shape.text if shape.has_text_frame else None,
                "placeholder_idx": (
                    shape.placeholder_format.idx
                    if shape.is_placeholder
                    else None
                ),
            }
            shapes.append(info)
        result.append({"slide_number": i, "shapes": shapes})

    return result


def pptx_replace_text(
    file_path: str,
    replacements: dict[str, str],
    output_path: str | None = None,
) -> dict:
    """Replace placeholder text across all slides in a .pptx file.

    Args:
        file_path: Path to the source .pptx file.
        replacements: Dict mapping placeholder strings to replacement values,
            e.g. {"{{title}}": "My Title", "{{date}}": "2026-02-13"}.
        output_path: Where to save the result. Defaults to overwriting the source.

    Returns a dict with counts of replacements made per placeholder.
    """
    from pptx import Presentation

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    counts = {k: 0 for k in replacements}

    def _replace_in_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                for old_text, new_text in replacements.items():
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        counts[old_text] += 1

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _replace_in_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_in_text_frame(cell.text_frame)
        # Also check speaker notes
        if slide.has_notes_slide:
            _replace_in_text_frame(slide.notes_slide.notes_text_frame)

    save_to = output_path or file_path
    prs.save(save_to)
    logger.info("Saved %s with %s replacements", save_to, sum(counts.values()))

    return {
        "output_path": save_to,
        "total_replacements": sum(counts.values()),
        "per_placeholder": counts,
    }


def pptx_get_slide_notes(file_path: str) -> list[dict]:
    """Read speaker notes from all slides.

    Returns a list of {slide_number, notes_text}.
    """
    from pptx import Presentation

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    result = []

    for i, slide in enumerate(prs.slides, 1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
        result.append({"slide_number": i, "notes_text": notes_text})

    return result


def pptx_duplicate_slide(
    file_path: str,
    slide_index: int,
    output_path: str | None = None,
) -> dict:
    """Duplicate a slide within a .pptx file.

    Args:
        file_path: Path to the source .pptx file.
        slide_index: 0-based index of the slide to duplicate.
        output_path: Where to save. Defaults to overwriting the source.

    Returns info about the duplicated slide.
    """
    from pptx import Presentation
    from copy import deepcopy
    from lxml import etree

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    slides = prs.slides

    if slide_index < 0 or slide_index >= len(slides):
        raise IndexError(
            f"Slide index {slide_index} out of range (0-{len(slides) - 1})"
        )

    # Get the slide to duplicate
    source_slide = slides[slide_index]

    # Copy the slide layout
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all elements from source to new slide
    for shape in source_slide.shapes:
        el = deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)

    # Remove the default placeholder shapes that add_slide creates
    # (they'd duplicate content if the layout has placeholders)
    for shape in list(new_slide.placeholders):
        sp = shape.element
        sp.getparent().remove(sp)

    # Re-add copied elements
    new_slide.shapes._spTree.clear()
    for child in deepcopy(source_slide.shapes._spTree):
        new_slide.shapes._spTree.append(child)

    save_to = output_path or file_path
    prs.save(save_to)

    return {
        "output_path": save_to,
        "source_slide_index": slide_index,
        "new_slide_index": len(slides) - 1,
        "total_slides": len(slides),
    }


# ── Google Slides API ─────────────────────────────────────────────────


def _get_slides_service():
    """Build an authenticated Google Slides API v1 service."""
    from roost.gmail.client import build_slides_service

    return build_slides_service()


def _get_drive_service():
    """Build an authenticated Google Drive API v3 service."""
    from googleapiclient.discovery import build
    from roost.gmail.client import _build_credentials

    creds = _build_credentials()
    return build("drive", "v3", credentials=creds)


def gslides_list_placeholders(presentation_id: str) -> list[dict]:
    """Read all slides and their text content to find placeholders.

    Returns a list of {slide_number, slide_id, elements: [{object_id, text}]}.
    """
    service = _get_slides_service()
    presentation = service.presentations().get(
        presentationId=presentation_id
    ).execute()

    result = []
    for i, slide in enumerate(presentation.get("slides", []), 1):
        elements = []
        for element in slide.get("pageElements", []):
            text = _extract_text_from_element(element)
            if text:
                elements.append({
                    "object_id": element.get("objectId"),
                    "text": text,
                })
        result.append({
            "slide_number": i,
            "slide_id": slide.get("objectId"),
            "elements": elements,
        })

    return result


def gslides_replace_text(
    presentation_id: str,
    replacements: dict[str, str],
) -> dict:
    """Replace placeholder text across all slides in a Google Slides presentation.

    Args:
        presentation_id: The Google Slides presentation ID.
        replacements: Dict mapping placeholder strings to replacement values.

    Returns counts of replacements made.
    """
    service = _get_slides_service()

    requests = []
    for old_text, new_text in replacements.items():
        requests.append({
            "replaceAllText": {
                "containsText": {
                    "text": old_text,
                    "matchCase": True,
                },
                "replaceText": new_text,
            }
        })

    if not requests:
        return {"total_replacements": 0, "per_placeholder": {}}

    response = service.presentations().batchUpdate(
        presentationId=presentation_id,
        body={"requests": requests},
    ).execute()

    # Parse replacement counts from response
    counts = {}
    replies = response.get("replies", [])
    placeholder_keys = list(replacements.keys())
    for idx, reply in enumerate(replies):
        replace_result = reply.get("replaceAllText", {})
        occurrences = replace_result.get("occurrencesChanged", 0)
        if idx < len(placeholder_keys):
            counts[placeholder_keys[idx]] = occurrences

    return {
        "presentation_id": presentation_id,
        "total_replacements": sum(counts.values()),
        "per_placeholder": counts,
    }


def gslides_create_from_template(
    template_id: str,
    name: str,
    replacements: dict[str, str] | None = None,
    folder_id: str | None = None,
) -> dict:
    """Copy a template presentation and optionally fill placeholders.

    Args:
        template_id: ID of the template presentation to copy.
        name: Name for the new presentation.
        replacements: Optional dict of placeholder → value replacements.
        folder_id: Optional Drive folder ID to place the copy in.

    Returns info about the created presentation.
    """
    drive_service = _get_drive_service()

    # Copy the template
    body = {"name": name}
    if folder_id:
        body["parents"] = [folder_id]

    copy = drive_service.files().copy(
        fileId=template_id,
        body=body,
    ).execute()

    new_id = copy["id"]
    logger.info("Created presentation '%s' (ID: %s) from template %s", name, new_id, template_id)

    result = {
        "presentation_id": new_id,
        "name": name,
        "template_id": template_id,
        "url": f"https://docs.google.com/presentation/d/{new_id}/edit",
    }

    # Apply replacements if provided
    if replacements:
        replace_result = gslides_replace_text(new_id, replacements)
        result["replacements"] = replace_result

    return result


def gslides_replace_image(
    presentation_id: str,
    placeholder_text: str,
    image_url: str,
) -> dict:
    """Replace shape placeholders containing specific text with an image.

    Args:
        presentation_id: The Google Slides presentation ID.
        placeholder_text: Text in the shape to match (e.g. "{{logo}}").
        image_url: Public URL of the image to insert.

    Returns count of shapes replaced.
    """
    service = _get_slides_service()

    requests = [{
        "replaceAllShapesWithImage": {
            "imageUrl": image_url,
            "imageReplaceMethod": "CENTER_INSIDE",
            "containsText": {
                "text": placeholder_text,
                "matchCase": True,
            },
        }
    }]

    response = service.presentations().batchUpdate(
        presentationId=presentation_id,
        body={"requests": requests},
    ).execute()

    replies = response.get("replies", [])
    occurrences = 0
    if replies:
        occurrences = replies[0].get("replaceAllShapesWithImage", {}).get(
            "occurrencesChanged", 0
        )

    return {
        "presentation_id": presentation_id,
        "placeholder_text": placeholder_text,
        "shapes_replaced": occurrences,
    }


# ── Helpers ───────────────────────────────────────────────────────────


def _extract_text_from_element(element: dict) -> str | None:
    """Extract all text from a Slides page element."""
    shape = element.get("shape", {})
    text_elements = shape.get("text", {}).get("textElements", [])
    if not text_elements:
        return None

    parts = []
    for te in text_elements:
        text_run = te.get("textRun", {})
        content = text_run.get("content", "")
        if content.strip():
            parts.append(content)

    return "".join(parts).strip() if parts else None
